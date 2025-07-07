from .document_processor import DocumentProcessor
from pptx import Presentation

class PptProcessor(DocumentProcessor):
    """
    PPT/PPTX 파일을 처리하는 프로세서. 프레젠테이션에서 텍스트를 추출한다.
    """
    extension = ('.ppt', '.pptx')

    def load(self, file_path: str) -> str:
        """
        PPT/PPTX 파일을 읽어 텍스트로 반환한다.

        Args:
            file_path (str): PPT/PPTX 파일 경로
        Returns:
            str: 추출된 텍스트
        """
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text_frame = getattr(shape, "text_frame", None)
                if text_frame and getattr(text_frame, "text", None):
                    text.append(text_frame.text)
        return "\n".join(text) 